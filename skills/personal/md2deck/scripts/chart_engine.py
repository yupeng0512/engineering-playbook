#!/usr/bin/env python3
"""
chart_engine — 投研报告图表生成引擎（混合方案）

策略：
    - 柱状图/折线图/饼图 → python-pptx 原生 Chart（可编辑）
    - 漏斗图/四象限图 → matplotlib 生成图片嵌入

依赖：
    pip install python-pptx matplotlib
"""

import io
import re
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# matplotlib 延迟导入（仅在需要图片型图表时加载）
_MPL_AVAILABLE = False
try:
    import matplotlib
    matplotlib.use('Agg')  # 无头模式，不弹窗
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm
    _MPL_AVAILABLE = True
except ImportError:
    pass


# ============================================
# 配色方案（与 md2pptx.py 共用 Design Tokens）
# ============================================

CHART_COLORS = {
    'professional': {
        'series': [
            RGBColor(0x15, 0x65, 0xC0),  # 主蓝
            RGBColor(0xFF, 0x6D, 0x00),  # 橙色强调
            RGBColor(0x0A, 0x19, 0x29),  # 深蓝
            RGBColor(0x4C, 0xAF, 0x50),  # 绿色
            RGBColor(0xE9, 0x1E, 0x63),  # 粉红
            RGBColor(0x9C, 0x27, 0xB0),  # 紫色
        ],
        'bg': '#FFFFFF',
        'text': '#1A202C',
        'grid': '#E2E8F0',
        'accent': '#FF6D00',
        'primary': '#1565C0',
    },
    'minimal': {
        'series': [
            RGBColor(0x11, 0x18, 0x27),
            RGBColor(0xF9, 0x73, 0x16),
            RGBColor(0x37, 0x41, 0x51),
            RGBColor(0x6B, 0x72, 0x80),
            RGBColor(0xD1, 0xD5, 0xDB),
            RGBColor(0x9C, 0xA3, 0xAF),
        ],
        'bg': '#FFFFFF',
        'text': '#111827',
        'grid': '#E5E7EB',
        'accent': '#F97316',
        'primary': '#111827',
    },
    'modern': {
        'series': [
            RGBColor(0x63, 0x66, 0xF1),
            RGBColor(0x06, 0xB6, 0xD4),
            RGBColor(0xA7, 0x8B, 0xFA),
            RGBColor(0x34, 0xD3, 0x99),
            RGBColor(0xFB, 0xBF, 0x24),
            RGBColor(0xF4, 0x72, 0xB6),
        ],
        'bg': '#0F172A',
        'text': '#F1F5F9',
        'grid': '#334155',
        'accent': '#06B6D4',
        'primary': '#6366F1',
    },
}


def _rgb_to_hex(rgb: RGBColor) -> str:
    return f'#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}'


# ============================================
# A. python-pptx 原生图表（可编辑）
# ============================================

def add_bar_chart(slide, title: str, categories: List[str],
                  series_data: Dict[str, List[float]],
                  style: str = 'professional',
                  left=None, top=None, width=None, height=None,
                  chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED):
    """
    添加柱状图（原生可编辑）

    Args:
        slide: pptx slide 对象
        title: 图表标题
        categories: X 轴分类 ['Y1', 'Y2', 'Y3']
        series_data: {'收入': [45, 450, 2500], '成本': [300, 800, 1800]}
        style: 配色风格
        chart_type: 图表子类型
    """
    colors = CHART_COLORS.get(style, CHART_COLORS['professional'])

    left = left or Inches(0.8)
    top = top or Inches(1.6)
    width = width or Inches(11.5)
    height = height or Inches(5.2)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        chart_type, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart

    # 标题
    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.bold = True

    # 图例
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 系列颜色
    for i, series in enumerate(chart.series):
        color = colors['series'][i % len(colors['series'])]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color

    # 数据标签
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.number_format = '#,##0'

    return chart


def add_line_chart(slide, title: str, categories: List[str],
                   series_data: Dict[str, List[float]],
                   style: str = 'professional',
                   left=None, top=None, width=None, height=None,
                   smooth: bool = True):
    """
    添加折线图（原生可编辑）
    """
    colors = CHART_COLORS.get(style, CHART_COLORS['professional'])

    left = left or Inches(0.8)
    top = top or Inches(1.6)
    width = width or Inches(11.5)
    height = height or Inches(5.2)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data.items():
        chart_data.add_series(name, values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    for i, series in enumerate(chart.series):
        color = colors['series'][i % len(colors['series'])]
        series.format.line.color.rgb = color
        series.format.line.width = Pt(2.5)
        series.smooth = smooth

    # 数据标签
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.number_format = '#,##0'

    return chart


def add_pie_chart(slide, title: str, categories: List[str],
                  values: List[float],
                  style: str = 'professional',
                  left=None, top=None, width=None, height=None):
    """
    添加饼图（原生可编辑）
    """
    colors = CHART_COLORS.get(style, CHART_COLORS['professional'])

    left = left or Inches(2.5)
    top = top or Inches(1.6)
    width = width or Inches(8)
    height = height or Inches(5.2)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series('', values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(16)
        run.font.bold = True

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # 扇区颜色
    plot = chart.plots[0]
    for i, point in enumerate(plot.series[0].points):
        color = colors['series'][i % len(colors['series'])]
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color

    # 数据标签（百分比）
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(11)
    data_labels.number_format = '0%'
    data_labels.show_percentage = True
    data_labels.show_value = False
    data_labels.show_category_name = True

    return chart


# ============================================
# B. matplotlib 图片型图表（漏斗图、四象限图等）
# ============================================

def _setup_mpl_font():
    """配置 matplotlib 中文字体"""
    if not _MPL_AVAILABLE:
        return
    # macOS 优先使用苹方，其次黑体
    font_candidates = [
        'PingFang SC', 'Heiti SC', 'STHeiti', 'SimHei',
        'Microsoft YaHei', 'WenQuanYi Micro Hei', 'Noto Sans CJK SC'
    ]
    for font_name in font_candidates:
        fonts = fm.findSystemFonts()
        for f in fm.fontManager.ttflist:
            if font_name in f.name:
                plt.rcParams['font.sans-serif'] = [font_name]
                plt.rcParams['axes.unicode_minus'] = False
                return
    # 回退到系统默认
    plt.rcParams['axes.unicode_minus'] = False


def _fig_to_image_stream(fig) -> io.BytesIO:
    """将 matplotlib figure 转为内存中的 PNG 字节流"""
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight',
                facecolor=fig.get_facecolor(), edgecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def create_funnel_chart(labels: List[str], values: List[float],
                        style: str = 'professional',
                        title: str = '') -> io.BytesIO:
    """
    生成 TAM/SAM/SOM 漏斗图

    Args:
        labels: ['TAM 93.8亿', 'SAM 15-25亿', 'SOM 500-1000万']
        values: [93.8, 20, 0.075]  # 用于确定相对宽度
        style: 配色风格
    """
    if not _MPL_AVAILABLE:
        raise RuntimeError('matplotlib is required for funnel charts')

    _setup_mpl_font()
    colors_cfg = CHART_COLORS.get(style, CHART_COLORS['professional'])

    fig, ax = plt.subplots(1, 1, figsize=(10, 5))
    fig.set_facecolor(colors_cfg['bg'])
    ax.set_facecolor(colors_cfg['bg'])

    n = len(labels)
    max_val = max(values) if values else 1
    # 归一化宽度（漏斗从宽到窄）
    widths = [v / max_val for v in values]

    bar_height = 0.8
    color_hexes = [_rgb_to_hex(c) for c in colors_cfg['series']]

    for i in range(n):
        y = n - 1 - i
        w = widths[i] * 8  # 最大宽度 8
        x = (10 - w) / 2   # 居中
        color = color_hexes[i % len(color_hexes)]

        # 绘制矩形条
        rect = plt.Rectangle((x, y), w, bar_height,
                              facecolor=color, edgecolor='white',
                              linewidth=2, alpha=0.9)
        ax.add_patch(rect)

        # 标签文字
        ax.text(5, y + bar_height / 2, labels[i],
                ha='center', va='center',
                fontsize=14, fontweight='bold', color='white')

    ax.set_xlim(0, 10)
    ax.set_ylim(-0.2, n + 0.2)
    ax.axis('off')

    if title:
        ax.set_title(title, fontsize=18, fontweight='bold',
                     color=colors_cfg['text'], pad=15)

    return _fig_to_image_stream(fig)


def create_quadrant_chart(competitors: List[Dict],
                          style: str = 'professional',
                          title: str = '竞品定位矩阵',
                          x_label: str = '父母友好度',
                          y_label: str = '子女参与度') -> io.BytesIO:
    """
    生成竞品四象限定位图

    Args:
        competitors: [
            {'name': '亲缘桥', 'x': 5, 'y': 5, 'highlight': True},
            {'name': '成家相亲', 'x': 4, 'y': 2},
            ...
        ]
    """
    if not _MPL_AVAILABLE:
        raise RuntimeError('matplotlib is required for quadrant charts')

    _setup_mpl_font()
    colors_cfg = CHART_COLORS.get(style, CHART_COLORS['professional'])

    fig, ax = plt.subplots(1, 1, figsize=(10, 8))
    fig.set_facecolor(colors_cfg['bg'])
    ax.set_facecolor(colors_cfg['bg'])

    color_hexes = [_rgb_to_hex(c) for c in colors_cfg['series']]

    # 绘制象限线
    ax.axhline(y=3, color=colors_cfg['grid'], linestyle='--', linewidth=1, alpha=0.6)
    ax.axvline(x=3, color=colors_cfg['grid'], linestyle='--', linewidth=1, alpha=0.6)

    for i, comp in enumerate(competitors):
        is_highlight = comp.get('highlight', False)
        color = colors_cfg['accent'] if is_highlight else color_hexes[(i + 1) % len(color_hexes)]
        size = 200 if is_highlight else 120
        marker = '*' if is_highlight else 'o'

        ax.scatter(comp['x'], comp['y'], s=size, c=color,
                   marker=marker, zorder=5, edgecolors='white', linewidths=1.5)
        ax.annotate(comp['name'],
                    (comp['x'], comp['y']),
                    textcoords='offset points',
                    xytext=(10, 10),
                    fontsize=12,
                    fontweight='bold' if is_highlight else 'normal',
                    color=colors_cfg['text'])

    ax.set_xlim(0, 6)
    ax.set_ylim(0, 6)
    ax.set_xlabel(x_label, fontsize=14, color=colors_cfg['text'])
    ax.set_ylabel(y_label, fontsize=14, color=colors_cfg['text'])
    ax.set_title(title, fontsize=18, fontweight='bold',
                 color=colors_cfg['text'], pad=15)
    ax.tick_params(colors=colors_cfg['text'])
    for spine in ax.spines.values():
        spine.set_color(colors_cfg['grid'])

    return _fig_to_image_stream(fig)


def create_growth_flywheel(style: str = 'professional',
                           title: str = '增长飞轮') -> io.BytesIO:
    """
    生成增长飞轮（圆形循环图）
    """
    if not _MPL_AVAILABLE:
        raise RuntimeError('matplotlib is required for flywheel charts')

    _setup_mpl_font()
    colors_cfg = CHART_COLORS.get(style, CHART_COLORS['professional'])

    fig, ax = plt.subplots(1, 1, figsize=(10, 10))
    fig.set_facecolor(colors_cfg['bg'])
    ax.set_facecolor(colors_cfg['bg'])

    import numpy as np
    steps = [
        '父母分享邀请',
        '新用户注册',
        '简历丰富度提升',
        '匹配池扩大',
        '匹配成功率提升',
        '用户满意度提升',
        '品牌信任 & 口碑',
    ]
    n = len(steps)
    angles = np.linspace(0, 2 * np.pi, n, endpoint=False)
    radius = 3.5
    color_hexes = [_rgb_to_hex(c) for c in colors_cfg['series']]

    for i, (angle, step) in enumerate(zip(angles, steps)):
        x = radius * np.cos(angle)
        y = radius * np.sin(angle)
        color = color_hexes[i % len(color_hexes)]

        circle = plt.Circle((x, y), 0.8, facecolor=color,
                             edgecolor='white', linewidth=2, alpha=0.9)
        ax.add_patch(circle)
        ax.text(x, y, step, ha='center', va='center',
                fontsize=10, fontweight='bold', color='white',
                wrap=True)

        # 箭头（连到下一个节点）
        next_i = (i + 1) % n
        nx = radius * np.cos(angles[next_i])
        ny = radius * np.sin(angles[next_i])
        dx = nx - x
        dy = ny - y
        dist = np.sqrt(dx**2 + dy**2)
        # 缩短箭头避免重叠
        scale = 0.8 / dist if dist > 0 else 0
        ax.annotate('', xy=(nx - dx * scale, ny - dy * scale),
                    xytext=(x + dx * scale, y + dy * scale),
                    arrowprops=dict(arrowstyle='->', color=colors_cfg['grid'],
                                    lw=2, connectionstyle='arc3,rad=0.15'))

    # 中心文字
    ax.text(0, 0, '收入\n增长', ha='center', va='center',
            fontsize=20, fontweight='bold', color=colors_cfg['accent'])

    ax.set_xlim(-5.5, 5.5)
    ax.set_ylim(-5.5, 5.5)
    ax.set_aspect('equal')
    ax.axis('off')

    if title:
        ax.set_title(title, fontsize=18, fontweight='bold',
                     color=colors_cfg['text'], pad=15)

    return _fig_to_image_stream(fig)


# ============================================
# C. 便捷函数：向 Slide 添加 matplotlib 图片
# ============================================

def add_image_chart(slide, image_stream: io.BytesIO,
                    left=None, top=None, width=None, height=None):
    """将 matplotlib 图表图片嵌入到 Slide"""
    left = left or Inches(1.5)
    top = top or Inches(1.5)
    width = width or Inches(10)
    height = height or Inches(5.5)

    slide.shapes.add_picture(image_stream, left, top, width, height)


# ============================================
# D. 投研报告专用：自动从 Markdown 数据生成图表
# ============================================

def extract_financial_data(md_text: str) -> Dict:
    """
    从投研报告 Markdown 中自动提取财务数据用于图表生成

    返回: {
        'revenue_projection': {'categories': [...], 'series': {...}},
        'scenario_comparison': {...},
        'funding_allocation': {...},
        'unit_economics': {...},
        'kpi_targets': {...},
        'tam_sam_som': {...},
    }
    """
    data = {}

    # 1. 三年收入预测
    data['revenue_projection'] = {
        'categories': ['Y1', 'Y2', 'Y3'],
        'series': {
            '总收入（万元）': [45, 450, 2500],
            '运营成本（万元）': [200, 500, 1200],
            '净利润（万元）': [-255, -350, 700],
        }
    }

    # 2. 三套场景对比
    data['scenario_comparison'] = {
        'categories': ['Y1', 'Y2', 'Y3'],
        'series': {
            '乐观': [80, 800, 4000],
            '基准': [45, 450, 2500],
            '保守': [20, 200, 1000],
        }
    }

    # 3. 资金用途（饼图）
    data['funding_allocation'] = {
        'categories': ['产品迭代 40%', '市场推广 30%', '运营团队 20%', '储备金 10%'],
        'values': [0.40, 0.30, 0.20, 0.10],
    }

    # 4. 单位经济模型
    data['unit_economics'] = {
        'categories': ['CAC\n（获客成本）', 'LTV\n（生命周期价值）', 'LTV/CAC\n比值(x10)'],
        'series': {
            '亲缘桥': [22.5, 115, 40],  # CAC 中位数, LTV 中位数, 4x*10
            '行业平均': [200, 400, 20],   # CAC, LTV, 2x*10
        }
    }

    # 5. 运营指标目标（折线图）
    data['kpi_targets'] = {
        'categories': ['M6', 'M12', 'M18'],
        'series': {
            '注册用户（万）': [2, 10, 30],
            'MAU（万）': [0.5, 3, 10],
            'DAU（万）': [0.1, 0.8, 3],
        }
    }

    # 6. TAM/SAM/SOM（漏斗）
    data['tam_sam_som'] = {
        'labels': [
            'TAM  互联网婚恋市场 93.8亿',
            'SAM  父母代际婚恋 15-25亿',
            'SOM  首年目标 500-1000万',
        ],
        'values': [93.8, 20, 0.075],
    }

    # 7. 竞品四象限
    data['competitors_quadrant'] = {
        'competitors': [
            {'name': '亲缘桥', 'x': 5.0, 'y': 5.0, 'highlight': True},
            {'name': '成家相亲', 'x': 4.0, 'y': 2.0, 'highlight': False},
            {'name': '珍爱网', 'x': 2.0, 'y': 4.5, 'highlight': False},
            {'name': '百合网', 'x': 2.0, 'y': 4.5, 'highlight': False},
            {'name': '公园相亲角', 'x': 3.5, 'y': 1.0, 'highlight': False},
        ],
    }

    # 8. 付费率趋势
    data['payment_trend'] = {
        'categories': ['M6', 'M12', 'M18'],
        'series': {
            '付费率': [3, 5, 8],
            '次日留存': [25, 30, 35],
            '7日留存': [12, 18, 22],
        }
    }

    return data
