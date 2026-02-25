#!/usr/bin/env python3
"""
md2pptx — 投研报告 Markdown → 可编辑 PowerPoint 转换器（v3.0）

特性:
    - 完整解析投研报告所有模块、子章节、表格、要点、引文
    - 每个子章节独立生成 Slide，确保数据完整呈现
    - 支持多风格：professional（默认）/ minimal / modern
    - 表格溢出自动拆分到多张 Slide
    - --compact 模式：只生成核心 Slide（12-18 张），适合路演
    - --charts 模式：自动生成图表（柱状图/折线图/饼图/漏斗图）
    - 混合图表引擎：原生 Chart（可编辑）+ matplotlib 图片（漏斗/四象限）

用法:
    python md2pptx.py --input report.md --output output.pptx [--style professional] [--compact] [--charts]

依赖:
    pip install python-pptx Pillow
    pip install matplotlib  # 可选，仅 --charts 需要
"""

import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ============================================
# Design Tokens — 三套风格
# ============================================

STYLES = {
    'professional': {
        'primary': RGBColor(0x15, 0x65, 0xC0),
        'primary_dark': RGBColor(0x0A, 0x19, 0x29),
        'primary_mid': RGBColor(0x1E, 0x3A, 0x5F),
        'accent': RGBColor(0xFF, 0x6D, 0x00),
        'accent_light': RGBColor(0xFF, 0xF3, 0xE0),
        'bg': RGBColor(0xFF, 0xFF, 0xFF),
        'bg_alt': RGBColor(0xF8, 0xFA, 0xFC),
        'text': RGBColor(0x1A, 0x20, 0x2C),
        'text_secondary': RGBColor(0x64, 0x74, 0x8B),
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'table_header': RGBColor(0x0A, 0x19, 0x29),
        'table_header_end': RGBColor(0x15, 0x65, 0xC0),
        'table_border': RGBColor(0xE2, 0xE8, 0xF0),
        'divider': RGBColor(0xFF, 0x6D, 0x00),
    },
    'minimal': {
        'primary': RGBColor(0x11, 0x18, 0x27),
        'primary_dark': RGBColor(0x03, 0x07, 0x12),
        'primary_mid': RGBColor(0x37, 0x41, 0x51),
        'accent': RGBColor(0xF9, 0x73, 0x16),
        'accent_light': RGBColor(0xFF, 0xF7, 0xED),
        'bg': RGBColor(0xFF, 0xFF, 0xFF),
        'bg_alt': RGBColor(0xFA, 0xFA, 0xFA),
        'text': RGBColor(0x11, 0x18, 0x27),
        'text_secondary': RGBColor(0x6B, 0x72, 0x80),
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'table_header': RGBColor(0x03, 0x07, 0x12),
        'table_header_end': RGBColor(0x37, 0x41, 0x51),
        'table_border': RGBColor(0xE5, 0xE7, 0xEB),
        'divider': RGBColor(0xF9, 0x73, 0x16),
    },
    'modern': {
        'primary': RGBColor(0x63, 0x66, 0xF1),
        'primary_dark': RGBColor(0x0F, 0x17, 0x2A),
        'primary_mid': RGBColor(0x31, 0x2E, 0x81),
        'accent': RGBColor(0x06, 0xB6, 0xD4),
        'accent_light': RGBColor(0xEC, 0xFD, 0xFF),
        'bg': RGBColor(0x0F, 0x17, 0x2A),
        'bg_alt': RGBColor(0x1E, 0x29, 0x3B),
        'text': RGBColor(0xF1, 0xF5, 0xF9),
        'text_secondary': RGBColor(0x94, 0xA3, 0xB8),
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'table_header': RGBColor(0x31, 0x2E, 0x81),
        'table_header_end': RGBColor(0x63, 0x66, 0xF1),
        'table_border': RGBColor(0x33, 0x41, 0x55),
        'divider': RGBColor(0x06, 0xB6, 0xD4),
    },
}

SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)


# ============================================
# Markdown Parser（增强版）
# ============================================

def parse_markdown(md_text):
    """
    解析 Markdown 为分层结构:
    - 顶层按 ## 分模块
    - 每个模块内按 ### 分子章节
    - 每个子章节独立收集 tables、bullets、quotes、text
    """
    modules = []
    current_module = None
    current_sub = None

    def _flush_sub():
        nonlocal current_sub
        if current_sub and current_module is not None:
            current_module['subsections'].append(current_sub)
        current_sub = None

    def _flush_module():
        nonlocal current_module
        _flush_sub()
        if current_module:
            modules.append(current_module)
        current_module = None

    def _new_sub(title=''):
        nonlocal current_sub
        _flush_sub()
        current_sub = {
            'title': title,
            'table_rows': [],
            'bullets': [],
            'quotes': [],
            'text_lines': [],
            'code_blocks': [],
        }

    in_code_block = False
    code_buf = []

    for line in md_text.split('\n'):
        # Code block toggle
        if line.strip().startswith('```'):
            if in_code_block:
                in_code_block = False
                if current_sub:
                    current_sub['code_blocks'].append('\n'.join(code_buf))
                code_buf = []
            else:
                in_code_block = True
                code_buf = []
            continue

        if in_code_block:
            code_buf.append(line)
            continue

        # H1 — 封面
        m_h1 = re.match(r'^#\s+(.+)', line)
        if m_h1 and not line.startswith('##'):
            _flush_module()
            current_module = {
                'title': m_h1.group(1).strip(),
                'subsections': [],
                'is_cover': True,
            }
            _new_sub()
            continue

        # H2 — 新模块
        m_h2 = re.match(r'^##\s+(.+)', line)
        if m_h2:
            _flush_module()
            current_module = {
                'title': m_h2.group(1).strip(),
                'subsections': [],
                'is_cover': False,
            }
            _new_sub()  # default subsection for content before first ###
            continue

        # H3 — 新子章节
        m_h3 = re.match(r'^###\s+(.+)', line)
        if m_h3:
            if current_module is None:
                continue
            _new_sub(m_h3.group(1).strip())
            continue

        if current_sub is None:
            continue

        # Table row
        if '|' in line and line.strip().startswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if cells and not all(re.match(r'^[-:]+$', c) for c in cells):
                current_sub['table_rows'].append(cells)
        # Bullet
        elif re.match(r'^\s*[-*]\s+', line):
            text = re.sub(r'^\s*[-*]\s+', '', line).strip()
            if text:
                current_sub['bullets'].append(text)
        # Numbered list
        elif re.match(r'^\s*\d+\.\s+', line):
            text = re.sub(r'^\s*\d+\.\s+', '', line).strip()
            if text:
                current_sub['bullets'].append(text)
        # Blockquote
        elif line.startswith('>'):
            text = line.lstrip('> ').strip()
            if text and not text.startswith('⚠️'):
                current_sub['quotes'].append(text)
        # Regular text
        elif line.strip():
            stripped = line.strip()
            # Skip --- dividers
            if stripped == '---':
                continue
            current_sub['text_lines'].append(stripped)

    _flush_module()
    return modules


def split_table(rows, max_rows=8):
    """将大表格拆分为多个小表格，每个不超过 max_rows 行数据"""
    if not rows or len(rows) < 2:
        return []
    header = rows[0]
    data = rows[1:]
    cols = len(header)
    # 过滤列数不匹配的行
    data = [r for r in data if len(r) == cols]

    tables = []
    for i in range(0, len(data), max_rows):
        chunk = data[i:i + max_rows]
        tables.append((header, chunk))
    return tables


# ============================================
# Slide Builders
# ============================================

def _set_font(run, size=18, bold=False, color=None, font_name=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if font_name:
        run.font.name = font_name


def _set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                  color=None, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """便捷方法：添加文本框并返回 text_frame"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    return tf


def _add_rich_text(paragraph, text, C):
    """解析 **加粗** 并用 accent 色标注"""
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            run = paragraph.add_run()
            run.text = part[2:-2]
            _set_font(run, size=18, bold=True, color=C['accent'])
        elif part.strip():
            run = paragraph.add_run()
            run.text = part
            _set_font(run, size=18, color=C['text'])


def _clean_md(text):
    """去除 Markdown 格式符号"""
    text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
    text = re.sub(r'\*([^*]+)\*', r'\1', text)
    text = re.sub(r'`([^`]+)`', r'\1', text)
    return text.strip()


# --- Slide types ---

def add_cover(prs, title, lines, C):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C['primary_dark'])

    # 底部装饰条
    _add_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), C['accent'])

    # 标题
    _add_text_box(slide, Inches(1), Inches(1.8), Inches(11.333), Inches(1.2),
                  _clean_md(title), size=52, bold=True, color=C['white'],
                  alignment=PP_ALIGN.CENTER)

    # 从内容提取副标题和融资信息
    subtitle = ''
    funding = ''
    for line in lines:
        clean = _clean_md(line)
        if not clean:
            continue
        if '融资' in clean or '天使轮' in clean:
            funding = clean
        elif '定位' in clean or '首个' in clean or '让' in clean:
            subtitle = clean
        elif not subtitle and len(clean) > 10:
            subtitle = clean

    if subtitle:
        _add_text_box(slide, Inches(1.5), Inches(3.2), Inches(10.333), Inches(1),
                      subtitle, size=26, color=RGBColor(0xCC, 0xDD, 0xEE),
                      alignment=PP_ALIGN.CENTER)

    if funding:
        # 融资信息用强调色胶囊
        _add_text_box(slide, Inches(3), Inches(4.6), Inches(7.333), Inches(0.8),
                      funding, size=24, bold=True, color=C['accent'],
                      alignment=PP_ALIGN.CENTER)

    # 底部信息
    _add_text_box(slide, Inches(1), Inches(6.2), Inches(11.333), Inches(0.5),
                  '微信小程序 · 云开发 · 2026', size=16,
                  color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)
    return slide


def add_section_slide(prs, title, C):
    """章节过渡页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C['bg_alt'] if C['bg'] == RGBColor(0xFF, 0xFF, 0xFF) else C['bg'])

    # 左侧装饰条
    _add_rect(slide, Inches(0), Inches(1.5), Inches(0.15), Inches(4.5), C['accent'])

    # 模块标题
    clean_title = _clean_md(title)
    # 提取模块编号（如"模块二：市场分析"）
    m = re.match(r'(模块[一二三四五六七八九十]+)[：:]\s*(.+)', clean_title)
    if m:
        _add_text_box(slide, Inches(1.2), Inches(2.2), Inches(10), Inches(0.8),
                      m.group(1), size=22, color=C['text_secondary'])
        _add_text_box(slide, Inches(1.2), Inches(3.0), Inches(10), Inches(1.5),
                      m.group(2), size=48, bold=True, color=C['primary_dark'])
    else:
        _add_text_box(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1.5),
                      clean_title, size=48, bold=True, color=C['primary_dark'])
    return slide


def add_title_slide(prs, title, C):
    """子章节标题"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
                  _clean_md(title), size=34, bold=True, color=C['primary_dark'])

    # 标题下划线
    _add_rect(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), C['accent'])

    return slide


def add_bullets_slide(prs, title, bullets, C, quotes=None):
    """标题 + 要点页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9),
                  _clean_md(title), size=34, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.2), Inches(2), Inches(0.06), C['accent'])

    # Bullets
    y = Inches(1.6)
    max_bullets = 7
    for i, bullet in enumerate(bullets[:max_bullets]):
        txBox = slide.shapes.add_textbox(Inches(1.2), y, Inches(11), Inches(0.55))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # 圆点标记
        dot = p.add_run()
        dot.text = '● '
        _set_font(dot, size=12, color=C['primary'])

        # 解析内容
        _add_rich_text(p, bullet, C)
        y += Inches(0.65)

    # 引文（如果有空间）
    if quotes and y < Inches(6.0):
        y += Inches(0.2)
        _add_rect(slide, Inches(1.2), y, Inches(0.08), Inches(0.5), C['accent'])
        quote_text = ' '.join(quotes)
        txBox_q = slide.shapes.add_textbox(Inches(1.6), y, Inches(10.5), Inches(0.8))
        tf_q = txBox_q.text_frame
        tf_q.word_wrap = True
        p_q = tf_q.paragraphs[0]
        run_q = p_q.add_run()
        run_q.text = _clean_md(quote_text)
        _set_font(run_q, size=16, color=C['text_secondary'])

    return slide


def add_table_slide(prs, title, header, rows, C, subtitle=''):
    """数据表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.8),
                  _clean_md(title), size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.1), Inches(1.5), Inches(0.05), C['accent'])

    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.4),
                      subtitle, size=16, color=C['text_secondary'])

    if not header or not rows:
        return slide

    cols = len(header)
    n_rows = len(rows) + 1

    table_top = Inches(1.5) if not subtitle else Inches(1.7)
    row_h = 0.42 if n_rows <= 7 else 0.36
    table_height = Inches(row_h * n_rows + 0.2)

    # 确保表格不超出 Slide 底部
    max_height = Inches(5.5) if not subtitle else Inches(5.3)
    if table_height > max_height:
        table_height = max_height

    table_shape = slide.shapes.add_table(
        n_rows, cols, Inches(0.8), table_top, Inches(11.5), table_height
    )
    table = table_shape.table

    # Header
    for j, h in enumerate(header):
        cell = table.cell(0, j)
        cell.text = _clean_md(h)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                _set_font(run, size=14, bold=True, color=C['white'])
        cell.fill.solid()
        cell.fill.fore_color.rgb = C['table_header']

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            if j >= cols:
                break
            cell = table.cell(i + 1, j)
            cell.text = _clean_md(val)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    _set_font(run, size=13, color=C['text'])
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C['bg_alt']

    return slide


def add_quote_slide(prs, quote_text, C):
    """引用/洞察页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C['bg_alt'] if C['bg'] == RGBColor(0xFF, 0xFF, 0xFF) else C['bg'])

    # 大引号
    _add_text_box(slide, Inches(2), Inches(1.5), Inches(1), Inches(1),
                  '\u201c', size=80, bold=True, color=C['accent'],
                  alignment=PP_ALIGN.LEFT)

    # 引文
    _add_text_box(slide, Inches(2), Inches(2.8), Inches(9.333), Inches(2.5),
                  _clean_md(quote_text), size=26, color=C['primary_dark'],
                  alignment=PP_ALIGN.CENTER)

    # 下引号
    _add_text_box(slide, Inches(10.333), Inches(5), Inches(1), Inches(1),
                  '\u201d', size=80, bold=True, color=C['accent'],
                  alignment=PP_ALIGN.RIGHT)

    return slide


def add_cta(prs, C):
    """结束页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, C['primary_dark'])
    _add_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), C['accent'])

    _add_text_box(slide, Inches(1), Inches(2.0), Inches(11.333), Inches(1.5),
                  '让每一个家庭都能找到好姻缘', size=44, bold=True,
                  color=C['white'], alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(3.6), Inches(11.333), Inches(1),
                  '亲缘桥 — 父母安心、子女放心的婚恋平台', size=24,
                  color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(5.2), Inches(11.333), Inches(0.5),
                  '联系方式：[待补充]', size=18,
                  color=RGBColor(0x88, 0x99, 0xAA), alignment=PP_ALIGN.CENTER)

    _add_text_box(slide, Inches(1), Inches(6.0), Inches(11.333), Inches(0.6),
                  '感谢您的时间', size=20, bold=True,
                  color=C['accent'], alignment=PP_ALIGN.CENTER)

    return slide


# ============================================
# Main Converter
# ============================================

def _is_core_module(title: str) -> bool:
    """判断是否为 compact 模式下的核心模块（严格）"""
    # compact 模式跳过：附录、数据来源索引、运营增长策略的细节
    skip_keywords = ['附录', '数据来源', 'PPT 路演大纲']
    if any(kw in title for kw in skip_keywords):
        return False

    core_keywords = [
        '执行摘要', 'Executive Summary',
        '市场分析', 'Market',
        '痛点', 'Problem', '机会', 'Opportunity',
        '产品', 'Product', '解决方案', 'Solution',
        '商业模式', 'Business Model',
        '竞品', 'Competitive',
        '财务', 'Financial',
        '团队', 'Team', '融资', 'Funding',
        '风险', 'Risk',
    ]
    return any(kw in title for kw in core_keywords)


def _is_core_subsection(title: str, module_title: str) -> bool:
    """在 compact 模式下，判断子章节是否核心（严格过滤）"""
    # 没有标题的子章节（模块默认内容）保留
    if not title:
        return True

    # 严格白名单：只保留对投资人最重要的子章节
    core_sub_keywords = [
        '一句话定位', '核心亮点', '融资需求',
        '市场规模', 'TAM', 'SAM',
        '核心痛点',
        '核心功能',
        '收入模型',
        '竞品矩阵', '竞争优势',
        '三年收入',
        '融资计划',
        '冷启动', '增长渠道',
    ]
    return any(kw in title for kw in core_sub_keywords)


def _add_chart_slides(prs, md_text: str, style: str, C: dict):
    """生成图表 Slide（使用混合引擎）"""
    try:
        from chart_engine import (
            extract_financial_data, add_bar_chart, add_line_chart,
            add_pie_chart, add_image_chart,
            create_funnel_chart, create_quadrant_chart,
            _MPL_AVAILABLE as mpl_ok,
        )
    except ImportError:
        # 尝试相对路径导入
        import importlib.util
        spec = importlib.util.spec_from_file_location(
            'chart_engine',
            str(Path(__file__).parent / 'chart_engine.py')
        )
        chart_engine = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(chart_engine)
        extract_financial_data = chart_engine.extract_financial_data
        add_bar_chart = chart_engine.add_bar_chart
        add_line_chart = chart_engine.add_line_chart
        add_pie_chart = chart_engine.add_pie_chart
        add_image_chart = chart_engine.add_image_chart
        create_funnel_chart = chart_engine.create_funnel_chart
        create_quadrant_chart = chart_engine.create_quadrant_chart
        mpl_ok = chart_engine._MPL_AVAILABLE

    data = extract_financial_data(md_text)

    # 图表章节过渡页
    add_section_slide(prs, '数据可视化', C)

    # 1. 三年收入预测（原生柱状图 — 可编辑）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                  '三年财务预测', size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
    rev = data['revenue_projection']
    add_bar_chart(slide, '收入 / 成本 / 净利润（万元）',
                  rev['categories'], rev['series'], style,
                  left=Inches(0.5), top=Inches(1.3),
                  width=Inches(12), height=Inches(5.8))

    # 2. 场景对比（原生柱状图）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                  '收入场景对比', size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
    sc = data['scenario_comparison']
    add_bar_chart(slide, '乐观 / 基准 / 保守 收入预测（万元）',
                  sc['categories'], sc['series'], style,
                  left=Inches(0.5), top=Inches(1.3),
                  width=Inches(12), height=Inches(5.8))

    # 3. 资金用途（原生饼图）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                  '资金用途分配', size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
    fa = data['funding_allocation']
    add_pie_chart(slide, '天使轮 300-500 万元资金用途',
                  fa['categories'], fa['values'], style,
                  left=Inches(2), top=Inches(1.3),
                  width=Inches(9), height=Inches(5.8))

    # 4. 运营指标趋势（原生折线图）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                  '关键运营指标趋势', size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
    kpi = data['kpi_targets']
    add_line_chart(slide, '注册用户 / MAU / DAU 增长（万）',
                   kpi['categories'], kpi['series'], style,
                   left=Inches(0.5), top=Inches(1.3),
                   width=Inches(12), height=Inches(5.8))

    # 5. 单位经济模型（原生柱状图）
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                  '单位经济模型', size=32, bold=True, color=C['primary_dark'])
    _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
    ue = data['unit_economics']
    add_bar_chart(slide, '亲缘桥 vs 行业平均',
                  ue['categories'], ue['series'], style,
                  left=Inches(0.5), top=Inches(1.3),
                  width=Inches(12), height=Inches(5.8))

    # 6-7. matplotlib 图片型图表（如果可用）
    if mpl_ok:
        # TAM/SAM/SOM 漏斗图
        tam = data['tam_sam_som']
        funnel_img = create_funnel_chart(
            tam['labels'], tam['values'], style, 'TAM / SAM / SOM 市场漏斗'
        )
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                      '市场规模漏斗', size=32, bold=True, color=C['primary_dark'])
        _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
        add_image_chart(slide, funnel_img,
                        left=Inches(1.5), top=Inches(1.5),
                        width=Inches(10), height=Inches(5.5))

        # 竞品四象限
        cq = data['competitors_quadrant']
        quad_img = create_quadrant_chart(
            cq['competitors'], style, '竞品定位矩阵',
            '父母友好度', '子女参与度'
        )
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8),
                      '竞品四象限定位', size=32, bold=True, color=C['primary_dark'])
        _add_rect(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(0.06), C['accent'])
        add_image_chart(slide, quad_img,
                        left=Inches(1.5), top=Inches(1.0),
                        width=Inches(10), height=Inches(6))

    chart_count = 7 if mpl_ok else 5
    print(f'   📊 图表: {chart_count} 张（原生可编辑 5 张' +
          (f' + matplotlib 图片 2 张）' if mpl_ok else '）'))


def convert(input_path, output_path, style='professional',
            template_path=None, compact=False, charts=False):
    md_text = Path(input_path).read_text(encoding='utf-8')
    modules = parse_markdown(md_text)
    C = STYLES.get(style, STYLES['professional'])

    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

    mode_label = 'compact' if compact else 'full'

    for module in modules:
        title = module['title']
        is_cover = module.get('is_cover', False)
        subs = module.get('subsections', [])

        # compact 模式：跳过非核心模块（但保留封面）
        if compact and not is_cover and not _is_core_module(title):
            continue

        # 封面
        if is_cover:
            all_lines = []
            for sub in subs:
                all_lines.extend(sub['text_lines'])
                all_lines.extend(sub['bullets'])
            add_cover(prs, title, all_lines, C)

            for sub in subs:
                if compact and sub['title'] and not _is_core_subsection(sub['title'], title):
                    continue
                if sub['title'] and sub['bullets']:
                    add_bullets_slide(prs, sub['title'], sub['bullets'], C, sub['quotes'])
                if sub['title'] and sub['table_rows']:
                    for header, data in split_table(sub['table_rows']):
                        add_table_slide(prs, sub['title'], header, data, C)
            continue

        # compact 模式：预检查 — 如果模块内所有子章节都会被过滤掉，
        # 则整个模块（含过渡页）都不生成，避免出现空模块
        if compact:
            has_visible_sub = False
            for sub in subs:
                if not sub['title'] or _is_core_subsection(sub['title'], title):
                    # 还需要检查该子章节是否真有内容
                    if (sub['table_rows'] or sub['bullets'] or
                            sub['quotes'] or sub['text_lines']):
                        has_visible_sub = True
                        break
            if not has_visible_sub:
                continue

        # 非封面模块 → 先加章节过渡页
        add_section_slide(prs, title, C)

        # 处理每个子章节
        for sub in subs:
            # compact 模式：跳过非核心子章节
            if compact and sub['title'] and not _is_core_subsection(sub['title'], title):
                continue

            sub_title = sub['title'] if sub['title'] else title
            has_table = bool(sub['table_rows'])
            has_bullets = bool(sub['bullets'])
            has_quotes = bool(sub['quotes'])
            has_text = bool(sub['text_lines'])

            if has_table:
                tables = split_table(sub['table_rows'])
                for idx, (header, data) in enumerate(tables):
                    suffix = f'（续 {idx+1}）' if idx > 0 else ''
                    add_table_slide(prs, sub_title + suffix, header, data, C)

                if has_bullets:
                    add_bullets_slide(prs, sub_title, sub['bullets'], C, sub['quotes'])
                elif has_quotes:
                    for q in sub['quotes']:
                        add_quote_slide(prs, q, C)

            elif has_bullets:
                chunk_size = 7
                for i in range(0, len(sub['bullets']), chunk_size):
                    chunk = sub['bullets'][i:i + chunk_size]
                    suffix = f'（续）' if i > 0 else ''
                    quotes_for_this = sub['quotes'] if i == 0 else None
                    add_bullets_slide(prs, sub_title + suffix, chunk, C, quotes_for_this)

            elif has_quotes:
                for q in sub['quotes']:
                    add_quote_slide(prs, q, C)

            elif has_text:
                meaningful = [l for l in sub['text_lines']
                              if len(l) > 5 and not l.startswith('#')]
                if meaningful:
                    add_bullets_slide(prs, sub_title, meaningful[:7], C)

    # 图表（如果启用）
    if charts:
        _add_chart_slides(prs, md_text, style, C)

    # 结束页
    add_cta(prs, C)

    prs.save(output_path)
    slide_count = len(prs.slides)
    print(f'✅ 已生成: {output_path}')
    print(f'   Slides: {slide_count} 张')
    print(f'   风格: {style}')
    print(f'   模式: {mode_label}')
    if charts:
        print(f'   图表: 已启用（混合引擎）')
    return slide_count


# ============================================
# CLI
# ============================================

def main():
    parser = argparse.ArgumentParser(
        description='md2pptx v3.0 — 投研报告 Markdown → 可编辑 PowerPoint'
    )
    parser.add_argument('--input', '-i', required=True, help='输入 Markdown 文件路径')
    parser.add_argument('--output', '-o', default='output.pptx', help='输出 .pptx 路径')
    parser.add_argument('--style', '-s', default='professional',
                        choices=['professional', 'minimal', 'modern'],
                        help='配色风格: professional（默认）/ minimal / modern')
    parser.add_argument('--template', '-t', default=None, help='PowerPoint 模板文件路径')
    parser.add_argument('--compact', '-c', action='store_true',
                        help='紧凑模式：只生成核心 Slide（12-18 张），适合路演')
    parser.add_argument('--charts', action='store_true',
                        help='启用图表生成（需要 matplotlib）')

    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f'❌ 输入文件不存在: {args.input}')
        sys.exit(1)

    convert(args.input, args.output, args.style, args.template,
            compact=args.compact, charts=args.charts)


if __name__ == '__main__':
    main()
